from pathlib import Path
from pypdf import PdfReader
from docx import Document
from pptx import Presentation


def load_text_from_file(file_path: Path) -> str:
    """
    Load text from TXT, PDF, DOCX, or PPTX files
    """
    suffix = file_path.suffix.lower()

    # ==========================
    # TXT
    # ==========================
    if suffix == ".txt":
        return file_path.read_text(encoding="utf-8", errors="ignore")

    # ==========================
    # PDF
    # ==========================
    if suffix == ".pdf":
        reader = PdfReader(str(file_path))
        text = ""

        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"

        return text.strip()

    # ==========================
    # DOCX
    # ==========================
    if suffix == ".docx":
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text.strip()

    # ==========================
    # PPTX
    # ==========================
    if suffix == ".pptx":
        prs = Presentation(file_path)
        text = ""

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

        return text.strip()

    raise ValueError(f"Unsupported file type: {suffix}")